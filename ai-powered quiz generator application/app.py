import os, json, re, time, requests, streamlit as st
from pptx import Presentation
from dotenv import load_dotenv
from io import BytesIO

load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "")
OPENROUTER_URL     = "https://openrouter.ai/api/v1/chat/completions"
MODEL              = "anthropic/claude-haiku-4-5"

DIFF_HINTS = {
    "Simple":  "Basic recall — definitions, key terms straight from the slides.",
    "Medium":  "Mix of recall + scenario-based reasoning.",
    "Complex": "Analytical, edge-cases & best-among-goods challenges.",
}
DIFF_PROMPTS = {
    "Simple":  "Generate straightforward factual recall questions focusing on definitions, key terms, and basic concepts directly stated in the slides.",
    "Medium":  "Generate balanced questions mixing factual recall and scenario-based reasoning.",
    "Complex": "Generate challenging analytical questions requiring deep understanding, edge cases, trade-offs, and novel application.",
}

st.set_page_config(page_title="AI Quiz Generator", page_icon="🎯", layout="wide", initial_sidebar_state="collapsed")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800;900&display=swap');

*, *::before, *::after { box-sizing:border-box; margin:0; padding:0; }

html, body {
    font-family:'Inter',sans-serif;
    height:100%; overflow:hidden;
}

[data-testid="stAppViewContainer"] {
    font-family:'Inter',sans-serif;
    background:#f8fafc;
    height:100vh;
    overflow:hidden;
}

[data-testid="stHeader"]      { display:none !important; }
[data-testid="stSidebar"]     { display:none !important; }
#MainMenu, footer             { visibility:hidden; }

[data-testid="stMain"] {
    height:100vh;
    overflow:hidden;
    padding:0 !important;
}

[data-testid="stMain"] > div {
    height:100vh;
    overflow:hidden;
    padding:0 !important;
}

[data-testid="block-container"] {
    padding: 0 !important;
    max-width: 100% !important;
    height: 100vh;
    overflow: hidden;
}

/* ── FULL PAGE WRAPPER ── */
.app-shell {
    display: flex;
    flex-direction: column;
    height: 100vh;
    background: #f8fafc;
    overflow: hidden;
}

/* ── TOPBAR ── */
.topbar {
    height: 58px;
    background: #fff;
    border-bottom: 1px solid #e2e8f0;
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 0 32px;
    flex-shrink: 0;
    box-shadow: 0 1px 4px rgba(0,0,0,0.06);
}
.tb-brand {
    display:flex; align-items:center; gap:10px;
    font-size:16px; font-weight:800; color:#0f172a;
}
.tb-logo {
    width:36px; height:36px; border-radius:10px;
    background:linear-gradient(135deg,#0f766e,#0891b2);
    display:flex; align-items:center; justify-content:center;
    font-size:18px;
    box-shadow:0 2px 8px rgba(15,118,110,0.35);
}
.tb-badge {
    background:#f0fdfa; color:#0f766e;
    border:1px solid #99f6e4;
    font-size:11px; font-weight:700;
    padding:5px 14px; border-radius:20px;
    letter-spacing:.04em;
}

/* ── MAIN AREA ── */
.main-area {
    flex: 1;
    display: flex;
    align-items: center;
    justify-content: center;
    overflow: hidden;
    padding: 24px;
}

/* ── CENTER CARD ── */
.center-card {
    background: #fff;
    border-radius: 24px;
    border: 1px solid #e2e8f0;
    box-shadow: 0 4px 24px rgba(0,0,0,0.07), 0 1px 4px rgba(0,0,0,0.04);
    width: 100%;
    max-width: 680px;
    padding: 44px 48px;
    overflow: hidden;
}

/* ── UPLOAD SCREEN ── */
.upload-zone {
    border: 2px dashed #99f6e4;
    border-radius: 16px;
    padding: 48px 24px;
    text-align: center;
    background: #f0fdfa;
    margin-bottom: 20px;
    transition: all .2s;
}
.upload-zone h3 { font-size:20px; font-weight:800; color:#0f172a; margin-bottom:6px; }
.upload-zone p  { font-size:13px; color:#94a3b8; }
.u-icon { font-size:48px; margin-bottom:14px; display:block; }

.parsed-ok {
    display:flex; align-items:center; gap:14px;
    background:#f0fdfa; border:1px solid #99f6e4;
    border-radius:12px; padding:14px 18px; margin:14px 0;
}
.parsed-ok strong { font-size:14px; color:#0f172a; display:block; margin-bottom:2px; }
.parsed-ok span   { font-size:12px; color:#64748b; }
.ok-badge { background:#0f766e; color:#fff; font-size:11px; font-weight:700; padding:4px 12px; border-radius:20px; }

/* ── SCREEN TITLE ── */
.screen-title { margin-bottom:28px; }
.screen-title h2 { font-size:26px; font-weight:800; color:#0f172a; margin-bottom:6px; letter-spacing:-.4px; }
.screen-title p  { font-size:14px; color:#64748b; }

/* ── SRC BAR ── */
.src-bar {
    background:#f0fdfa; border:1px solid #99f6e4;
    border-radius:10px; padding:10px 16px;
    font-size:13px; color:#0f766e; font-weight:600;
    margin-bottom:22px; display:flex; align-items:center; gap:8px;
}
.src-bar em { color:#64748b; font-style:normal; font-weight:400; }

/* ── DIFF BUTTONS ── */
.diff-row { display:flex; gap:10px; margin:10px 0 6px; }
.diff-btn {
    flex:1; padding:16px 10px; border-radius:12px;
    border:2px solid #e2e8f0; background:#f8fafc;
    text-align:center; cursor:pointer; transition:all .2s;
}
.diff-btn.on { border-color:#0f766e; background:#f0fdfa; box-shadow:0 0 0 3px #ccfbf1; }
.diff-btn .de { font-size:22px; margin-bottom:5px; display:block; }
.diff-btn .dn { font-size:13px; font-weight:700; color:#374151; }
.diff-btn.on .dn { color:#0f766e; }
.hint { background:#f0f9ff; border:1px solid #bae6fd; border-radius:8px; padding:9px 13px; font-size:12px; color:#0369a1; margin-bottom:18px; }

/* ── QUIZ HEADER ── */
.quiz-hdr {
    display:flex; align-items:center; justify-content:space-between;
    margin-bottom:18px;
}
.q-progress-label { font-size:13px; font-weight:700; color:#0f766e; }
.q-timer {
    background:#fff7ed; border:1px solid #fed7aa;
    color:#c2410c; font-size:13px; font-weight:700;
    padding:6px 14px; border-radius:10px;
    display:flex; align-items:center; gap:5px;
}

.pbar-wrap { background:#e2e8f0; border-radius:6px; height:6px; margin-bottom:26px; overflow:hidden; }
.pbar-fill { height:100%; background:linear-gradient(90deg,#0f766e,#0891b2); border-radius:6px; transition:width .4s; }

.q-tag { display:inline-flex; align-items:center; gap:6px; background:linear-gradient(135deg,#0f766e,#0891b2); color:#fff; font-size:11px; font-weight:700; padding:4px 12px; border-radius:6px; margin-bottom:12px; }
.q-body { font-size:17px; font-weight:700; color:#0f172a; line-height:1.55; margin-bottom:22px; letter-spacing:-.2px; }

/* ── OPTION BUTTONS (key override) ── */
.opt-unsel, .opt-sel {
    display:flex; align-items:center; gap:14px;
    width:100%; padding:14px 18px;
    border-radius:12px; cursor:pointer;
    margin-bottom:10px; transition:all .15s;
    font-size:14px; font-weight:500;
    border:2px solid #e2e8f0;
    background:#f8fafc; color:#0f172a;
}
.opt-sel {
    border-color:#0f766e !important;
    background:#f0fdfa !important;
    color:#0f766e !important; font-weight:700 !important;
    box-shadow:0 0 0 3px #ccfbf1 !important;
}
.opt-key {
    width:30px; height:30px; border-radius:8px;
    display:flex; align-items:center; justify-content:center;
    font-size:13px; font-weight:800; flex-shrink:0;
    background:#e2e8f0; color:#475569;
}
.opt-sel .opt-key { background:#0f766e; color:#fff; }

/* Nav row */
.nav-row { display:flex; justify-content:space-between; align-items:center; margin-top:18px; }

/* ── SCORE SCREEN ── */
.score-wrap {
    background:linear-gradient(135deg,#0f766e,#0891b2);
    border-radius:18px; padding:32px 28px; text-align:center;
    margin-bottom:22px; position:relative; overflow:hidden;
}
.score-wrap::before { content:''; position:absolute; top:-50px; right:-50px; width:180px; height:180px; border-radius:50%; background:rgba(255,255,255,0.07); }
.score-num { font-size:68px; font-weight:900; color:#fff; line-height:1; position:relative; z-index:1; }
.score-pct { font-size:20px; color:rgba(255,255,255,.7); font-weight:600; margin-bottom:6px; position:relative; z-index:1; }
.score-msg { font-size:16px; font-weight:700; color:#fff; margin-bottom:18px; position:relative; z-index:1; }
.chips { display:flex; justify-content:center; gap:10px; flex-wrap:wrap; position:relative; z-index:1; }
.chip { padding:6px 16px; border-radius:30px; font-size:12px; font-weight:700; }
.cg { background:rgba(16,185,129,.22); color:#6ee7b7; border:1px solid rgba(16,185,129,.4); }
.cr { background:rgba(239,68,68,.18); color:#fca5a5; border:1px solid rgba(239,68,68,.35); }
.cb { background:rgba(255,255,255,.13); color:rgba(255,255,255,.85); border:1px solid rgba(255,255,255,.25); }

.rev-list { max-height:300px; overflow-y:auto; padding-right:4px; }
.rev-list::-webkit-scrollbar { width:4px; }
.rev-list::-webkit-scrollbar-track { background:#f1f5f9; border-radius:4px; }
.rev-list::-webkit-scrollbar-thumb { background:#cbd5e1; border-radius:4px; }

.ritem { border-radius:12px; padding:14px 16px; margin-bottom:9px; }
.ritem.ok  { background:#f0fdf4; border:1px solid #bbf7d0; }
.ritem.bad { background:#fef2f2; border:1px solid #fecaca; }
.ritem-hd  { display:flex; align-items:flex-start; gap:8px; margin-bottom:4px; }
.ritem-tl  { font-size:13px; font-weight:700; color:#0f172a; }
.ritem-q   { font-size:12px; color:#64748b; margin:3px 0 4px 23px; }
.ans-ok  { font-size:12px; font-weight:600; color:#16a34a; margin-left:23px; }
.ans-bad { font-size:12px; font-weight:600; color:#dc2626; margin-left:23px; }
.ai-exp  { background:#f0f9ff; border:1px solid #bae6fd; border-radius:8px; padding:9px 12px; font-size:12px; color:#0369a1; margin:7px 0 0; display:flex; gap:7px; }

/* ── STREAMLIT OVERRIDES ── */
/* Hide all default widget decorations */
[data-testid="stFileUploaderDropzone"] {
    background:#f0fdfa !important;
    border:2px dashed #99f6e4 !important;
    border-radius:14px !important;
}
[data-testid="stFileUploaderDropzone"] p     { color:#64748b !important; }
[data-testid="stFileUploaderDropzone"] small { color:#9ca3af !important; }

div[data-testid="stSlider"] label { color:#0f172a !important; font-weight:700 !important; font-size:13px !important; }
[data-testid="stSlider"] [role="slider"] { background:#0f766e !important; box-shadow:0 0 0 4px #ccfbf1 !important; }

.stButton > button {
    border-radius:12px !important; font-weight:700 !important;
    font-size:14px !important; padding:13px 20px !important;
    transition:all .2s !important; width:100% !important;
}
.stButton > button[kind="primary"] {
    background:linear-gradient(135deg,#0f766e,#0891b2) !important;
    color:#fff !important; border:none !important;
    box-shadow:0 4px 16px rgba(15,118,110,0.4) !important;
}
.stButton > button[kind="primary"]:hover {
    box-shadow:0 6px 24px rgba(15,118,110,0.55) !important;
    transform:translateY(-2px) !important;
}
.stButton > button[kind="secondary"] {
    background:#fff !important; border:2px solid #e2e8f0 !important; color:#374151 !important;
}
.stButton > button[kind="secondary"]:hover {
    border-color:#0f766e !important; color:#0f766e !important; background:#f0fdfa !important;
}
.stButton > button:disabled { opacity:.35 !important; transform:none !important; box-shadow:none !important; }
[data-testid="stSpinner"] > div { border-top-color:#0f766e !important; }
.stCaption,[data-testid="stCaptionContainer"] p { color:#94a3b8 !important; font-size:11px !important; }

/* Remove streamlit default padding everywhere */
.element-container, .stMarkdown { margin:0 !important; }
</style>
""", unsafe_allow_html=True)

# ── SESSION STATE ─────────────────────────────────────────────────────────────
def init():
    for k, v in {
        "step":"upload","slides":[],"filename":"","slide_count":0,"word_count":0,
        "num_questions":10,"difficulty":"Medium","questions":[],"user_answers":{},
        "current_q":0,"start_time":None,"elapsed":0,"explanations":{},"rk":0,
    }.items():
        if k not in st.session_state: st.session_state[k] = v
init()

# ── BACKEND ───────────────────────────────────────────────────────────────────
def extract_pptx(b):
    prs = Presentation(BytesIO(b)); slides, words = [], 0
    for i, slide in enumerate(prs.slides, 1):
        txt = [p.text.strip() for s in slide.shapes if s.has_text_frame for p in s.text_frame.paragraphs if p.text.strip()]
        c = " ".join(txt)
        if c: slides.append({"slide":i,"text":c}); words += len(c.split())
    return slides, words

def llm(prompt):
    r = requests.post(OPENROUTER_URL, headers={"Authorization":f"Bearer {OPENROUTER_API_KEY}","Content-Type":"application/json","HTTP-Referer":"http://localhost:8501","X-Title":"AI Quiz Generator"},
        json={"model":MODEL,"messages":[{"role":"user","content":prompt}],"temperature":.7,"max_tokens":4096}, timeout=120)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"]

def gen_mcqs(slides, n, diff):
    txt = "\n\n".join(f"[Slide {s['slide']}]: {s['text']}" for s in slides)
    raw = llm(f"Expert quiz designer. Generate exactly {n} MCQs.\nDIFFICULTY: {diff} — {DIFF_PROMPTS[diff]}\nSLIDES:\n{txt}\nRules: 4 options A-D, one correct, plausible distractors.\nReturn ONLY valid JSON array:\n[{{\"id\":1,\"question\":\"...\",\"options\":{{\"A\":\"...\",\"B\":\"...\",\"C\":\"...\",\"D\":\"...\"}},\"correct\":\"B\",\"topic\":\"label\"}}]")
    return json.loads(re.sub(r"```(?:json)?","",raw).strip().strip("`").strip())

def gen_exp(questions, answers):
    wrong = [q for q in questions if answers.get(str(q["id"])) and answers[str(q["id"])] != q["correct"]]
    if not wrong: return {}
    block = "\n\n".join(f"Q{w['id']}: {w['question']}\nStudent: {answers[str(w['id'])]} — {w['options'].get(answers[str(w['id'])],'')} \nCorrect: {w['correct']} — {w['options'][w['correct']]}" for w in wrong)
    raw = llm(f"Write 1-2 sentences per wrong answer explaining why it's wrong and why correct is right.\n{block}\nReturn ONLY JSON: {{\"1\":\"...\"}}")
    return json.loads(re.sub(r"```(?:json)?","",raw).strip().strip("`").strip())

# ── TOP BAR ───────────────────────────────────────────────────────────────────
badges = {"upload":"Step 1 / 3","config":"Step 2 / 3","quiz":"In Progress","results":"Complete ✓"}
st.markdown(f"""
<div class="topbar">
  <div class="tb-brand"><div class="tb-logo">🎯</div>AI Quiz Generator</div>
  <span class="tb-badge">{badges.get(st.session_state.step,'')}</span>
</div>
""", unsafe_allow_html=True)

# ── CENTER CARD WRAPPER ───────────────────────────────────────────────────────
_, col, _ = st.columns([1, 2.8, 1])

with col:

    # ══════════════════════════════════════════════════════════════════════════
    # SCREEN 1 — UPLOAD
    # ══════════════════════════════════════════════════════════════════════════
    if st.session_state.step == "upload":
        st.markdown("""
        <div class="screen-title">
          <h2>Upload Presentation</h2>
          <p>Import your .pptx — AI will extract all content automatically</p>
        </div>
        <div class="upload-zone">
          <span class="u-icon">📂</span>
          <h3>Drag & drop your .pptx here</h3>
          <p>.pptx only &nbsp;·&nbsp; max 25 MB</p>
        </div>
        """, unsafe_allow_html=True)

        up = st.file_uploader("Upload", type=["pptx"], label_visibility="collapsed")
        if up:
            if up.size > 25*1024*1024:
                st.error("File exceeds 25 MB.")
            else:
                with st.spinner("Parsing slides…"):
                    try:
                        slides, wc = extract_pptx(up.read())
                        if not slides: st.error("No text found.")
                        else:
                            st.session_state.update(slides=slides, filename=up.name, slide_count=len(slides), word_count=wc)
                            st.markdown(f"""
                            <div class="parsed-ok">
                              <span style="font-size:22px">📄</span>
                              <div style="flex:1"><strong>{up.name}</strong><span>{len(slides)} slides · {wc:,} words</span></div>
                              <span class="ok-badge">✓ Ready</span>
                            </div>
                            """, unsafe_allow_html=True)
                    except Exception as e:
                        st.error(f"Failed: {e}")

        st.write("")
        if st.button("Continue →", type="primary", disabled=not bool(st.session_state.slides)):
            st.session_state.step = "config"; st.rerun()

    # ══════════════════════════════════════════════════════════════════════════
    # SCREEN 2 — CONFIGURE
    # ══════════════════════════════════════════════════════════════════════════
    elif st.session_state.step == "config":
        st.markdown("""
        <div class="screen-title">
          <h2>Configure Your Quiz</h2>
          <p>Choose difficulty and question count</p>
        </div>
        """, unsafe_allow_html=True)

        st.markdown(f'<div class="src-bar">📎 <strong>{st.session_state.filename}</strong> <em>· {st.session_state.slide_count} slides · {st.session_state.word_count:,} words</em></div>', unsafe_allow_html=True)

        st.session_state.num_questions = st.slider("Number of questions", 5, 30, st.session_state.num_questions)

        st.markdown("<br>**Difficulty Level**", unsafe_allow_html=True)
        c1,c2,c3 = st.columns(3)
        for col2, (name,emoji) in zip([c1,c2,c3],[("Simple","🟢"),("Medium","🟡"),("Complex","🔴")]):
            with col2:
                on = "on" if st.session_state.difficulty==name else ""
                st.markdown(f'<div class="diff-btn {on}"><span class="de">{emoji}</span><span class="dn">{name}</span></div>', unsafe_allow_html=True)
                if st.button(name, key=f"d_{name}", type="primary" if on else "secondary"):
                    st.session_state.difficulty=name; st.rerun()

        st.markdown(f'<div class="hint" style="margin-top:10px">💡 {DIFF_HINTS[st.session_state.difficulty]}</div>', unsafe_allow_html=True)
        st.write("")
        if st.button(f"⚡  Generate {st.session_state.num_questions} Questions", type="primary"):
            with st.spinner("AI crafting your quiz…"):
                try:
                    qs = gen_mcqs(st.session_state.slides, st.session_state.num_questions, st.session_state.difficulty)
                    st.session_state.update(questions=qs, user_answers={}, current_q=0, start_time=time.time(), step="quiz", rk=st.session_state.rk+1)
                    st.rerun()
                except json.JSONDecodeError: st.error("AI returned malformed JSON — retry.")
                except Exception as e: st.error(f"Error: {e}")

    # ══════════════════════════════════════════════════════════════════════════
    # SCREEN 3 — QUIZ (one question, full card)
    # ══════════════════════════════════════════════════════════════════════════
    elif st.session_state.step == "quiz":
        qs      = st.session_state.questions
        total   = len(qs)
        i       = st.session_state.current_q
        q       = qs[i]
        qid     = str(q["id"])
        elapsed = int(time.time()-(st.session_state.start_time or time.time()))
        m,s     = divmod(elapsed,60)
        pct     = int(((i)/total)*100)
        sel     = st.session_state.user_answers.get(qid)

        st.markdown(f"""
        <div class="quiz-hdr">
          <span class="q-progress-label">Question {i+1} of {total}</span>
          <span class="q-timer">⏱ {m:02d}:{s:02d}</span>
        </div>
        <div class="pbar-wrap"><div class="pbar-fill" style="width:{pct}%"></div></div>
        <span class="q-tag">Q{i+1} &nbsp;·&nbsp; {q.get('topic','')}</span>
        <div class="q-body">{q['question']}</div>
        """, unsafe_allow_html=True)

        # Option buttons — styled via CSS, one button per option
        for key,val in q["options"].items():
            css = "opt-sel" if sel==key else "opt-unsel"
            st.markdown(f"""
            <div class="{css}" style="margin-bottom:0">
              <span class="opt-key">{key}</span>
              <span>{val}</span>
            </div>
            """, unsafe_allow_html=True)
            if st.button(val, key=f"o_{i}_{key}_{st.session_state.rk}", use_container_width=True, type="primary" if sel==key else "secondary"):
                st.session_state.user_answers[qid]=key; st.rerun()

        # Prev / Next
        cp,_,cn = st.columns([1,.2,1])
        with cp:
            if i>0:
                if st.button("← Previous", type="secondary"):
                    st.session_state.current_q-=1; st.rerun()
        with cn:
            lbl = "Submit Quiz ✓" if i==total-1 else "Next →"
            if st.button(lbl, type="primary"):
                if i<total-1:
                    st.session_state.current_q+=1; st.rerun()
                else:
                    st.session_state.elapsed=int(time.time()-st.session_state.start_time)
                    with st.spinner("Generating AI feedback…"):
                        try: st.session_state.explanations=gen_exp(qs, st.session_state.user_answers)
                        except: st.session_state.explanations={}
                    st.session_state.step="results"; st.rerun()

    # ══════════════════════════════════════════════════════════════════════════
    # SCREEN 4 — RESULTS
    # ══════════════════════════════════════════════════════════════════════════
    elif st.session_state.step == "results":
        qs      = st.session_state.questions
        total   = len(qs)
        correct = sum(1 for q in qs if st.session_state.user_answers.get(str(q["id"]))==q["correct"])
        wrong   = total-correct
        pct     = round((correct/total)*100) if total else 0
        m,s     = divmod(st.session_state.elapsed,60)
        hl      = ("Perfect! 🎉" if pct==100 else "Excellent! 🌟" if pct>=90 else "Great job! 💪" if pct>=80 else "Good effort 📚" if pct>=60 else "Keep going! 🚀")

        st.markdown(f"""
        <div class="score-wrap">
          <div class="score-num">{correct}/{total}</div>
          <div class="score-pct">{pct}%</div>
          <div class="score-msg">{hl}</div>
          <div class="chips">
            <span class="chip cg">✓ {correct} Correct</span>
            <span class="chip cr">✗ {wrong} Wrong</span>
            <span class="chip cb">⏱ {m:02d}:{s:02d}</span>
          </div>
        </div>
        <div style="font-size:11px;font-weight:700;letter-spacing:.08em;color:#94a3b8;text-transform:uppercase;margin-bottom:10px">Review & AI Feedback</div>
        <div class="rev-list">
        """, unsafe_allow_html=True)

        for idx,q in enumerate(qs):
            qid=str(q["id"]); ua=st.session_state.user_answers.get(qid)
            ok=ua==q["correct"]; exp=st.session_state.explanations.get(qid,"")
            if ok:   ans=f'<div class="ans-ok">✓ {ua}: {q["options"].get(ua,"")} — correct</div>'
            elif ua: ans=f'<div class="ans-bad">✗ You: {ua}: {q["options"].get(ua,"")} · Correct: {q["correct"]}: {q["options"][q["correct"]]}</div>'
            else:    ans=f'<div class="ans-bad">✗ Not answered · Correct: {q["correct"]}: {q["options"][q["correct"]]}</div>'
            ai=f'<div class="ai-exp"><span>🤖</span><span>{exp}</span></div>' if (not ok and exp) else ""
            st.markdown(f'<div class="ritem {"ok" if ok else "bad"}"><div class="ritem-hd"><span>{"✅" if ok else "❌"}</span><span class="ritem-tl">Q{idx+1} · {q.get("topic","")}</span></div><div class="ritem-q">{q["question"]}</div>{ans}{ai}</div>', unsafe_allow_html=True)

        st.markdown('</div>', unsafe_allow_html=True)
        st.write("")
        c1,c2=st.columns(2)
        with c1:
            if st.button("↺ Retake Quiz", type="primary"):
                st.session_state.update(user_answers={},current_q=0,start_time=time.time(),explanations={},step="quiz",rk=st.session_state.rk+1)
                st.rerun()
        with c2:
            if st.button("↑ New PPT", type="secondary"):
                for k in list(st.session_state.keys()): del st.session_state[k]
                st.rerun()
